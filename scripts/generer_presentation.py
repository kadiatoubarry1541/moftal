#!/usr/bin/env python3
"""
Script pour générer la présentation PowerPoint du projet "1"
Nécessite: pip install python-pptx
"""

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor
except ImportError:
    print("❌ Erreur: python-pptx n'est pas installé.")
    print("📦 Installez-le avec: pip install python-pptx")
    exit(1)

# Créer une nouvelle présentation
prs = Presentation()

# Définir les dimensions (16:9)
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Fonction pour ajouter un titre
def add_title(slide, title_text):
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

# Fonction pour ajouter du contenu
def add_content(slide, content_list):
    left = Inches(0.5)
    top = Inches(2)
    width = Inches(9)
    height = Inches(5)
    
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True
    
    for i, item in enumerate(content_list):
        p = text_frame.paragraphs[i] if i < len(text_frame.paragraphs) else text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(18)
        p.space_after = Pt(12)
        if i == 0:
            p.font.bold = True

# SLIDE 1 : PAGE DE COUVERTURE
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title1 = slide1.shapes.title
title1.text = "1"
subtitle = slide1.placeholders[1]
subtitle.text = "Plateforme Communautaire Guinéenne"

# SLIDE 2 : PROBLÉMATIQUE
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
add_title(slide2, "Le Problème")
add_content(slide2, [
    "• Les Guinéens de la diaspora sont dispersés",
    "• Difficile de maintenir les liens familiaux",
    "• Manque de plateforme centralisée",
    "• Perte de l'histoire et de la généalogie",
    "• Accès limité aux services de l'État"
])

# SLIDE 3 : SOLUTION
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
add_title(slide3, "Notre Solution")
add_content(slide3, [
    "🔗 Connecte les familles et préserve la généalogie",
    "🏘️ Organise la communauté par résidence et région",
    "🎯 Facilite les activités sociales",
    "🏥 Gère la santé communautaire",
    "📋 Simplifie les services de l'État",
    "💼 Facilite les échanges commerciaux",
    "🎓 Offre l'éducation et les formations",
    "🕌 Gère la foi et les dons (zakat)"
])

# SLIDE 4 : FONCTIONNALITÉS
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
add_title(slide4, "Fonctionnalités Clés")
add_content(slide4, [
    "1. Gestion Familiale - Arbre généalogique visuel",
    "2. Communauté - Groupes par résidence, région",
    "3. Services de l'État - Rendez-vous et documents",
    "4. Échanges - Vente et achat de produits",
    "5. Santé & Éducation - Formations et suivi"
])

# SLIDE 5 : TECHNOLOGIES
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
add_title(slide5, "Technologies Utilisées")
add_content(slide5, [
    "Frontend: React 18 + TypeScript + Vite",
    "Backend: Node.js + Express + PostgreSQL",
    "5 langues: fr, en, ar, man, pul",
    "Paiement mobile intégré",
    "Système de badges et logos"
])

# SLIDE 6 : INNOVATION
slide6 = prs.slides.add_slide(prs.slide_layouts[1])
add_title(slide6, "Ce qui nous rend unique")
add_content(slide6, [
    "✅ Système NumeroH unique",
    "✅ Interface multilingue (5 langues)",
    "✅ Plateforme tout-en-un",
    "✅ Généalogie complète",
    "✅ Paiement intégré",
    "✅ Sécurisé avec JWT"
])

# SLIDE 7 : IMPACT
slide7 = prs.slides.add_slide(prs.slide_layouts[1])
add_title(slide7, "Impact sur la Communauté")
add_content(slide7, [
    "🌍 Rapprochement de la diaspora",
    "👨‍👩‍👧‍👦 Préservation de l'histoire familiale",
    "💼 Facilitation des échanges commerciaux",
    "🏛️ Simplification des démarches administratives",
    "🎓 Accès facilité aux formations",
    "🤝 Gestion transparente des dons"
])

# SLIDE 8 : MODÈLE ÉCONOMIQUE
slide8 = prs.slides.add_slide(prs.slide_layouts[1])
add_title(slide8, "Modèle Économique")
add_content(slide8, [
    "💵 Dons et contributions volontaires",
    "🕌 Zakat et dons religieux",
    "🤝 Partenariats avec organisations",
    "🏛️ Subventions (si applicable)",
    "",
    "Partenaires: Opérateurs mobiles, Organisations guinéennes"
])

# SLIDE 9 : DÉMO
slide9 = prs.slides.add_slide(prs.slide_layouts[1])
add_title(slide9, "Aperçu de la Plateforme")
add_content(slide9, [
    "[Insérer captures d'écran:]",
    "1. Page d'accueil avec logo '1'",
    "2. Arbre généalogique",
    "3. Services de l'État",
    "4. Page des échanges"
])

# SLIDE 10 : ROADMAP
slide10 = prs.slides.add_slide(prs.slide_layouts[1])
add_title(slide10, "Prochaines Étapes")
add_content(slide10, [
    "Court terme:",
    "  ✅ Développement des fonctionnalités de base",
    "  🔄 Tests utilisateurs",
    "",
    "Moyen terme:",
    "  📢 Lancement beta",
    "  🤝 Partenariats",
    "",
    "Long terme:",
    "  📱 Application mobile",
    "  🌍 Expansion"
])

# SLIDE 11 : ÉQUIPE
slide11 = prs.slides.add_slide(prs.slide_layouts[1])
add_title(slide11, "Notre Équipe")
add_content(slide11, [
    "[Votre nom] - Fondateur & Développeur",
    "",
    "Vision:",
    "Créer une plateforme qui unit et sert",
    "la communauté guinéenne"
])

# SLIDE 12 : MERCI
slide12 = prs.slides.add_slide(prs.slide_layouts[0])
title12 = slide12.shapes.title
title12.text = "Merci pour votre attention !"
subtitle12 = slide12.placeholders[1]
subtitle12.text = "Questions ?\n\n1 - Plateforme Communautaire Guinéenne"

# Sauvegarder
output_file = "Presentation_Projet_1.pptx"
prs.save(output_file)
print(f"✅ Présentation créée avec succès: {output_file}")
print(f"📊 Nombre de slides: {len(prs.slides)}")

